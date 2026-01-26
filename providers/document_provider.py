#!/usr/bin/env python3
# -*- coding: utf-8 -*-
"""
Document Provider voor Baarn Raadsinformatie Server.
Downloads en verwerkt documenten van Notubiz.

Storage modes:
- KEEP_PDF_FILES=false (default): Download PDF, extract text, delete PDF
  → Alleen tekst in database, minimale disk space
- KEEP_PDF_FILES=true: Download PDF, extract text, keep PDF
  → Tekst in database + originele PDFs op disk
- STORE_FILES_IN_DB=true (default): Bestand bytes opslaan in database (BLOB)
"""

import base64
import io
import os
import requests
import tempfile
from pathlib import Path
from typing import Dict, List, Optional, Tuple
from urllib.parse import urlparse, unquote

from core.config import Config
from core.database import Database, get_database
from shared.logging_config import get_logger, LogContext

logger = get_logger('document-provider')

# PDF text extraction
try:
    import pdfplumber
    PDF_SUPPORT = True
    # Suppress pdfminer warnings about malformed PDF color values and fonts
    import logging
    logging.getLogger('pdfminer').setLevel(logging.ERROR)
except ImportError:
    PDF_SUPPORT = False
    logger.warning('pdfplumber not installed - PDF text extraction disabled')

# PDF image extraction (optional)
try:
    import fitz  # PyMuPDF
    PDF_IMAGE_SUPPORT = True
except ImportError:
    PDF_IMAGE_SUPPORT = False
    logger.warning('PyMuPDF not installed - PDF image extraction disabled')

# DOCX support
try:
    import docx
    DOCX_SUPPORT = True
except ImportError:
    DOCX_SUPPORT = False
    logger.warning('python-docx not installed - DOCX extraction disabled')

# PPTX support
try:
    from pptx import Presentation
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    PPTX_SUPPORT = True
except ImportError:
    PPTX_SUPPORT = False
    logger.warning('python-pptx not installed - PPTX extraction disabled')

# XLSX support
try:
    import openpyxl
    XLSX_SUPPORT = True
except ImportError:
    XLSX_SUPPORT = False
    logger.warning('openpyxl not installed - XLSX extraction disabled')

# OCR support (for image text extraction)
try:
    import pytesseract
    from PIL import Image
    OCR_SUPPORT = True
except ImportError:
    OCR_SUPPORT = False
    logger.warning('pytesseract/Pillow not installed - OCR disabled')

# Image deduplication (perceptual hashing)
try:
    import imagehash
    IMAGE_HASH_SUPPORT = True
except ImportError:
    IMAGE_HASH_SUPPORT = False
    logger.warning('imagehash not installed - image deduplication disabled')


class DocumentProvider:
    """Provider voor document downloads en text extractie."""

    def __init__(self, db: Database = None):
        """Initialize document provider."""
        self.db = db or get_database()
        self.documents_dir = Config.DOCUMENTS_DIR
        self.images_dir = Config.DATA_DIR / 'images'
        self.shared_images_dir = self.images_dir / 'shared'  # Deduplicated images
        self.keep_files = getattr(Config, 'KEEP_PDF_FILES', False)
        self.store_files_in_db = getattr(Config, 'STORE_FILES_IN_DB', False)  # Default to False now
        self.max_file_size_bytes = getattr(Config, 'MAX_FILE_SIZE_MB', 25) * 1024 * 1024

        # Ensure directories exist
        self.images_dir.mkdir(parents=True, exist_ok=True)
        self.shared_images_dir.mkdir(parents=True, exist_ok=True)

        if self.keep_files:
            self.documents_dir.mkdir(parents=True, exist_ok=True)

        mode = "keeping PDFs" if self.keep_files else "text-only (PDFs deleted after extraction)"
        dedup = "enabled" if IMAGE_HASH_SUPPORT else "disabled"
        logger.info(f'DocumentProvider initialized: {mode}, image deduplication: {dedup}')

    def _compute_image_hash(self, image_bytes: bytes) -> Optional[str]:
        """Compute perceptual hash of an image for deduplication."""
        if not IMAGE_HASH_SUPPORT:
            return None
        try:
            from PIL import Image
            with Image.open(io.BytesIO(image_bytes)) as img:
                # Use perceptual hash (pHash) - robust to minor changes
                phash = imagehash.phash(img)
                return str(phash)
        except Exception as e:
            logger.debug(f'Failed to compute image hash: {e}')
            return None

    def _save_image_to_filesystem(self, document_id: int, image_index: int,
                                   image_bytes: bytes, ext: str) -> Dict[str, any]:
        """
        Save an extracted image to the filesystem with deduplication.

        If the image already exists (same perceptual hash), references the
        existing shared image instead of creating a duplicate.

        Args:
            document_id: ID of the source document
            image_index: Index of the image within the document
            image_bytes: Raw image bytes
            ext: File extension (png, jpg, etc.)

        Returns:
            Dict with image metadata including file_path, image_hash, unique_image_id
        """
        # Get image dimensions
        width, height = None, None
        try:
            from PIL import Image
            with Image.open(io.BytesIO(image_bytes)) as img:
                width, height = img.size
        except Exception:
            pass

        # Compute perceptual hash for deduplication
        image_hash = self._compute_image_hash(image_bytes)

        # Check for existing image with same hash
        if image_hash:
            existing = self.db.find_unique_image_by_hash(image_hash)
            if existing:
                # Image already exists, reference it
                self.db.increment_unique_image_reference(existing['id'])
                logger.debug(f'Found duplicate image (hash={image_hash[:8]}...), referencing existing')
                return {
                    'index': image_index,
                    'mime_type': f'image/{ext}',
                    'file_path': existing['file_path'],
                    'image_hash': image_hash,
                    'unique_image_id': existing['id'],
                    'file_size': existing['file_size'],
                    'width': existing['width'],
                    'height': existing['height'],
                    'is_duplicate': True
                }

        # New unique image - save to shared directory
        if image_hash:
            # Use hash as filename for shared images
            filename = f'{image_hash}.{ext}'
            file_path = self.shared_images_dir / filename
        else:
            # Fallback: save to document-specific directory
            doc_image_dir = self.images_dir / f'doc_{document_id}'
            doc_image_dir.mkdir(parents=True, exist_ok=True)
            filename = f'{image_index:03d}.{ext}'
            file_path = doc_image_dir / filename

        # Write image file
        file_path.write_bytes(image_bytes)

        # Register in unique_images if we have a hash
        unique_image_id = None
        if image_hash:
            unique_image_id = self.db.add_unique_image(
                image_hash=image_hash,
                file_path=str(file_path),
                mime_type=f'image/{ext}',
                width=width,
                height=height,
                file_size=len(image_bytes)
            )
            logger.debug(f'Added new unique image (hash={image_hash[:8]}...)')

        return {
            'index': image_index,
            'mime_type': f'image/{ext}',
            'file_path': str(file_path),
            'image_hash': image_hash,
            'unique_image_id': unique_image_id,
            'file_size': len(image_bytes),
            'width': width,
            'height': height,
            'is_duplicate': False
        }

    def _cleanup_document_images(self, document_id: int):
        """
        Remove image records and files for a document.
        Handles deduplicated images by decrementing reference counts.
        Only deletes shared images when reference count reaches 0.
        """
        # Get paths to delete from database (handles reference counting)
        paths_to_delete = self.db.clear_document_images(document_id)

        # Delete the actual files
        for file_path in paths_to_delete:
            try:
                path = Path(file_path)
                if path.exists():
                    path.unlink()
                    logger.debug(f'Deleted image file: {file_path}')
            except Exception as e:
                logger.warning(f'Failed to delete image file {file_path}: {e}')

        # Also clean up document-specific directory if empty
        doc_image_dir = self.images_dir / f'doc_{document_id}'
        if doc_image_dir.exists():
            try:
                # Only remove if empty
                if not any(doc_image_dir.iterdir()):
                    doc_image_dir.rmdir()
            except Exception:
                pass

    def download_pending_documents(self, limit: int = None) -> Tuple[int, int]:
        """
        Download alle pending documents.

        Args:
            limit: Maximum number of documents to download

        Returns:
            Tuple of (successful, failed) downloads
        """
        pending = self.db.get_documents_pending_download()
        if limit:
            pending = pending[:limit]

        with LogContext(logger, 'download_pending', count=len(pending)):
            success = 0
            failed = 0

            for doc in pending:
                if self.download_document(doc['id']):
                    success += 1
                else:
                    failed += 1

            logger.info(f'Downloaded {success}/{len(pending)} documents')
            return success, failed

    def create_document_from_base64(
        self,
        title: str,
        filename: str,
        mime_type: str,
        file_base64: str,
        source_url: str = None
    ) -> int:
        """Create a document from base64 file data and extract content."""
        file_bytes = base64.b64decode(file_base64)
        if len(file_bytes) > self.max_file_size_bytes:
            raise ValueError('File too large for DB storage')

        document_id = self.db.upsert_document(
            title=title,
            url=source_url,
            filename=filename,
            mime_type=mime_type,
            file_size=len(file_bytes)
        )

        if self.store_files_in_db:
            self.db.update_document_file_blob(document_id, file_bytes, storage_mode='db')

        ext = Path(filename).suffix or '.bin'
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=ext)
        temp_path = Path(temp_file.name)
        temp_file.close()
        try:
            temp_path.write_bytes(file_bytes)
            # Clean up old images BEFORE extracting new ones
            self._cleanup_document_images(document_id)
            self.db.clear_document_images(document_id)

            text_content, images = self._extract_content_from_bytes(temp_path, file_bytes, document_id)
            if text_content:
                self.db.update_document_content(document_id, text_content)
            if images:
                self.db.add_document_images(document_id, images)
        finally:
            if temp_path.exists():
                temp_path.unlink()

        self.db.update_document_status(document_id, 'stored')
        return document_id

    def download_document(self, document_id: int) -> bool:
        """
        Download a single document, extract text, optionally delete PDF.

        Storage mode:
        - KEEP_PDF_FILES=false: Download to temp, extract text, delete file
        - KEEP_PDF_FILES=true: Download to documents dir, extract text, keep file

        Args:
            document_id: Database ID of document

        Returns:
            True if successful
        """
        doc = self.db.get_document(document_id)
        if not doc:
            logger.error(f'Document not found: {document_id}')
            return False

        url = doc.get('url')
        if not url:
            logger.warning(f'No URL for document {document_id}')
            self.db.update_document_status(document_id, 'no_url')
            return False

        try:
            # Determine where to save
            if self.keep_files:
                local_path = self._generate_local_path(doc)
                temp_file = None
            else:
                # Use temp file that will be deleted
                suffix = Path(urlparse(url).path).suffix or '.pdf'
                temp_file = tempfile.NamedTemporaryFile(delete=False, suffix=suffix)
                local_path = Path(temp_file.name)
                temp_file.close()

            # Download file
            logger.debug(f'Downloading: {url}')
            response = requests.get(url, timeout=60, stream=True)
            response.raise_for_status()

            # Save file
            with open(local_path, 'wb') as f:
                for chunk in response.iter_content(chunk_size=8192):
                    f.write(chunk)

            file_size = local_path.stat().st_size
            logger.debug(f'Downloaded {file_size} bytes to {local_path}')

            # Read bytes for storage/extraction
            file_bytes = local_path.read_bytes()

            if self.store_files_in_db:
                if len(file_bytes) > self.max_file_size_bytes:
                    logger.warning(f'Document {document_id}: file too large for DB storage ({len(file_bytes)} bytes)')
                    self.db.update_document_status(document_id, 'file_too_large')
                    if local_path.exists() and (not self.keep_files or self.store_files_in_db):
                        local_path.unlink()
                    return False
                self.db.update_document_file_blob(document_id, file_bytes, storage_mode='db')

            # Clean up old images BEFORE extracting new ones
            self._cleanup_document_images(document_id)
            self.db.clear_document_images(document_id)

            # Extract text and images (images saved to filesystem during extraction)
            text_content, images = self._extract_content_from_bytes(local_path, file_bytes, document_id)
            if images:
                self.db.add_document_images(document_id, images)

            # Update database status
            if self.store_files_in_db:
                self.db.update_document_status(document_id, 'stored')
            elif self.keep_files:
                self.db.update_document_status(document_id, 'downloaded', str(local_path))
            else:
                self.db.update_document_status(document_id, 'text_extracted', None)

            if text_content:
                self.db.update_document_content(document_id, text_content)
                logger.info(f'Document {document_id}: extracted {len(text_content)} chars')
            else:
                logger.warning(f'Document {document_id}: no text extracted')

            # Cleanup temp file if not keeping or if stored in DB
            if local_path.exists() and (not self.keep_files or self.store_files_in_db):
                local_path.unlink()
                logger.debug(f'Deleted temp file: {local_path}')

            return True

        except requests.exceptions.RequestException as e:
            logger.error(f'Download failed for document {document_id}: {e}')
            self.db.update_document_status(document_id, 'download_failed')
            return False
        except Exception as e:
            logger.error(f'Error processing document {document_id}: {e}')
            self.db.update_document_status(document_id, 'error')
            # Cleanup on error
            if 'local_path' in locals() and local_path.exists() and (not self.keep_files or self.store_files_in_db):
                local_path.unlink()
            return False

    def _generate_local_path(self, doc: Dict) -> Path:
        """Generate local file path for document."""
        # Get filename from document or URL
        filename = doc.get('filename')
        if not filename and doc.get('url'):
            url_path = urlparse(doc['url']).path
            filename = unquote(os.path.basename(url_path))

        if not filename:
            # Generate from title
            title = doc.get('title', 'document')
            # Sanitize title
            safe_title = "".join(c for c in title if c.isalnum() or c in (' ', '-', '_'))[:50]
            filename = f"{safe_title}.pdf"

        # Add document ID to ensure uniqueness
        doc_id = doc.get('id', 0)
        name, ext = os.path.splitext(filename)
        unique_filename = f"{doc_id}_{name}{ext}"

        return self.documents_dir / unique_filename

    def _extract_text_from_pdf_bytes(self, file_bytes: bytes) -> Optional[str]:
        """Extract text content from PDF bytes."""
        if not PDF_SUPPORT:
            logger.warning('PDF extraction not available - pdfplumber not installed')
            return None
        try:
            text_parts = []
            with pdfplumber.open(io.BytesIO(file_bytes)) as pdf:
                for page in pdf.pages:
                    page_text = page.extract_text()
                    if page_text:
                        text_parts.append(page_text)
            return '\n\n'.join(text_parts) if text_parts else None
        except Exception as e:
            logger.error(f'PDF text extraction failed: {e}')
            return None

    def _extract_pdf_images(self, file_bytes: bytes, document_id: int = None) -> List[Dict[str, any]]:
        """Extract images from PDF bytes using PyMuPDF and save to filesystem."""
        if not PDF_IMAGE_SUPPORT:
            return []
        images = []
        try:
            doc = fitz.open(stream=file_bytes, filetype='pdf')
            index = 0
            for page in doc:
                for img in page.get_images(full=True):
                    xref = img[0]
                    base = doc.extract_image(xref)
                    image_bytes = base.get('image')
                    if not image_bytes:
                        continue
                    ext = base.get('ext', 'png')

                    if document_id:
                        # Save to filesystem
                        image_meta = self._save_image_to_filesystem(
                            document_id, index, image_bytes, ext
                        )
                        images.append(image_meta)
                    else:
                        # Fallback to base64 (for compatibility)
                        images.append({
                            'index': index,
                            'mime_type': f'image/{ext}',
                            'data_base64': base64.b64encode(image_bytes).decode('ascii')
                        })
                    index += 1
            return images
        except Exception as e:
            logger.warning(f'PDF image extraction failed: {e}')
            return []

    def _extract_docx_content(self, file_bytes: bytes, document_id: int = None) -> Tuple[Optional[str], List[Dict[str, any]]]:
        """Extract text and images from DOCX bytes."""
        if not DOCX_SUPPORT:
            return None, []
        try:
            doc = docx.Document(io.BytesIO(file_bytes))
            text_parts = [p.text for p in doc.paragraphs if p.text]
            images = []
            index = 0
            for part in doc.part.related_parts.values():
                if getattr(part, 'content_type', '').startswith('image/'):
                    image_bytes = part.blob
                    # Determine extension from content type
                    ext = part.content_type.split('/')[-1]
                    if ext == 'jpeg':
                        ext = 'jpg'

                    if document_id:
                        image_meta = self._save_image_to_filesystem(
                            document_id, index, image_bytes, ext
                        )
                        images.append(image_meta)
                    else:
                        images.append({
                            'index': index,
                            'mime_type': part.content_type,
                            'data_base64': base64.b64encode(image_bytes).decode('ascii')
                        })
                    index += 1
            return '\n'.join(text_parts) if text_parts else None, images
        except Exception as e:
            logger.warning(f'DOCX extraction failed: {e}')
            return None, []

    def _extract_pptx_content(self, file_bytes: bytes, document_id: int = None) -> Tuple[Optional[str], List[Dict[str, any]]]:
        """Extract text and images from PPTX bytes."""
        if not PPTX_SUPPORT:
            return None, []
        try:
            presentation = Presentation(io.BytesIO(file_bytes))
            text_parts = []
            images = []
            index = 0
            for slide in presentation.slides:
                for shape in slide.shapes:
                    if hasattr(shape, 'text') and shape.text:
                        text_parts.append(shape.text)
                    if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        image = shape.image
                        ext = image.content_type.split('/')[-1]
                        if ext == 'jpeg':
                            ext = 'jpg'

                        if document_id:
                            image_meta = self._save_image_to_filesystem(
                                document_id, index, image.blob, ext
                            )
                            images.append(image_meta)
                        else:
                            images.append({
                                'index': index,
                                'mime_type': image.content_type,
                                'data_base64': base64.b64encode(image.blob).decode('ascii')
                            })
                        index += 1
            return '\n'.join(text_parts) if text_parts else None, images
        except Exception as e:
            logger.warning(f'PPTX extraction failed: {e}')
            return None, []

    def _extract_xlsx_content(self, file_bytes: bytes) -> Tuple[Optional[str], List[Dict[str, str]]]:
        """Extract text and images from XLSX bytes."""
        if not XLSX_SUPPORT:
            return None, []
        try:
            workbook = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True)
            text_parts = []
            images = []
            index = 0
            for sheet in workbook.worksheets:
                for row in sheet.iter_rows(values_only=True):
                    row_values = [str(cell) for cell in row if cell is not None]
                    if row_values:
                        text_parts.append('\t'.join(row_values))
                for img in getattr(sheet, '_images', []):
                    image_bytes = None
                    if hasattr(img, '_data') and callable(img._data):
                        image_bytes = img._data()
                    elif hasattr(img, '_data'):
                        image_bytes = img._data
                    if image_bytes:
                        mime_type = 'image/png'
                        images.append({
                            'index': index,
                            'mime_type': mime_type,
                            'data_base64': base64.b64encode(image_bytes).decode('ascii')
                        })
                        index += 1
            return '\n'.join(text_parts) if text_parts else None, images
        except Exception as e:
            logger.warning(f'XLSX extraction failed: {e}')
            return None, []

    def _extract_content_from_bytes(
        self,
        file_path: Path,
        file_bytes: bytes,
        document_id: int = None
    ) -> Tuple[Optional[str], Optional[List[Dict[str, str]]]]:
        """Extract text and images from supported file types.

        If document_id is provided, images are saved to filesystem.
        """
        ext = file_path.suffix.lower()
        if ext == '.pdf':
            text = self._extract_text_from_pdf_bytes(file_bytes)
            images = self._extract_pdf_images(file_bytes, document_id)
            return text, images
        if ext == '.docx':
            return self._extract_docx_content(file_bytes, document_id)
        if ext == '.pptx':
            return self._extract_pptx_content(file_bytes, document_id)
        if ext == '.xlsx':
            return self._extract_xlsx_content(file_bytes)
        return None, []

    def extract_text(self, document_id: int) -> Optional[str]:
        """
        Extract text content from an already-downloaded document.
        Only works if KEEP_PDF_FILES=true.

        Args:
            document_id: Database ID of document

        Returns:
            Extracted text or None
        """
        doc = self.db.get_document(document_id)
        if not doc:
            return None

        local_path = doc.get('local_path')
        if not local_path or not Path(local_path).exists():
            logger.warning(f'Document file not found: {document_id}')
            return None

        file_path = Path(local_path)
        file_bytes = file_path.read_bytes()

        # Clean up old images BEFORE extracting new ones
        self._cleanup_document_images(document_id)
        self.db.clear_document_images(document_id)

        full_text, images = self._extract_content_from_bytes(file_path, file_bytes, document_id)

        if full_text:
            self.db.update_document_content(document_id, full_text)
            logger.info(f'Extracted text from document {document_id}: {len(full_text)} chars')
        if images:
            self.db.add_document_images(document_id, images)

        return full_text

    def extract_all_text(self) -> Tuple[int, int]:
        """
        Extract text from all downloaded documents without extracted text.

        Returns:
            Tuple of (successful, failed) extractions
        """
        # Find documents needing extraction
        docs = self.db.get_documents()
        need_extraction = [
            d for d in docs
            if d.get('download_status') == 'downloaded'
            and not d.get('text_extracted')
            and d.get('local_path')
        ]

        success = 0
        failed = 0

        with LogContext(logger, 'extract_all_text', count=len(need_extraction)):
            for doc in need_extraction:
                if self.extract_text(doc['id']):
                    success += 1
                else:
                    failed += 1

        return success, failed

    # ==================== OCR ====================

    def ocr_image(self, image_path: str, lang: str = 'nld') -> Optional[str]:
        """
        Extract text from an image using OCR.

        Args:
            image_path: Path to the image file
            lang: Tesseract language code (nld=Dutch, eng=English)

        Returns:
            Extracted text or None if failed
        """
        if not OCR_SUPPORT:
            logger.warning('OCR not available - pytesseract not installed')
            return None

        try:
            image = Image.open(image_path)
            # Use Dutch + English for best results on Dutch government docs
            text = pytesseract.image_to_string(image, lang=f'{lang}+eng')
            return text.strip() if text.strip() else None
        except Exception as e:
            logger.warning(f'OCR failed for {image_path}: {e}')
            return None

    def process_pending_ocr(self, limit: int = 100) -> Tuple[int, int]:
        """
        Process images that need OCR.
        Handles both unique (deduplicated) and document-specific images.

        Args:
            limit: Maximum number of images to process

        Returns:
            Tuple of (successful, failed) OCR operations
        """
        if not OCR_SUPPORT:
            logger.warning('OCR not available - pytesseract not installed')
            return 0, 0

        success = 0
        failed = 0

        # First, process unique (deduplicated) images - OCR is stored once
        unique_pending = self.db.get_unique_images_pending_ocr(limit)
        with LogContext(logger, 'ocr_unique_images', count=len(unique_pending)):
            for image in unique_pending:
                file_path = image.get('file_path')
                if not file_path or not Path(file_path).exists():
                    self.db.update_unique_image_ocr(image['id'], None, 'file_missing')
                    failed += 1
                    continue

                ocr_text = self.ocr_image(file_path)
                if ocr_text:
                    self.db.update_unique_image_ocr(image['id'], ocr_text, 'completed')
                    success += 1
                    logger.debug(f"OCR completed for unique image {image['id']}: {len(ocr_text)} chars")
                else:
                    self.db.update_unique_image_ocr(image['id'], '', 'no_text')
                    success += 1

        # Then process non-deduplicated images (fallback for images without hash)
        remaining_limit = max(0, limit - len(unique_pending))
        if remaining_limit > 0:
            pending = self.db.get_images_pending_ocr(remaining_limit)
            # Filter out images that reference unique_images (already processed above)
            pending = [img for img in pending if not img.get('unique_image_id')]

            with LogContext(logger, 'ocr_processing', count=len(pending)):
                for image in pending:
                    file_path = image.get('file_path')
                    if not file_path or not Path(file_path).exists():
                        self.db.update_image_ocr(image['id'], None, 'file_missing')
                        failed += 1
                        continue

                    ocr_text = self.ocr_image(file_path)
                    if ocr_text:
                        self.db.update_image_ocr(image['id'], ocr_text, 'completed')
                        success += 1
                        logger.debug(f"OCR completed for image {image['id']}: {len(ocr_text)} chars")
                    else:
                        self.db.update_image_ocr(image['id'], '', 'no_text')
                        success += 1

        logger.info(f'OCR processing: {success} successful, {failed} failed')
        return success, failed

    def get_deduplication_stats(self) -> Dict:
        """
        Get image deduplication statistics.

        Returns:
            Dict with deduplication stats
        """
        stats = self.db.get_deduplication_stats()

        # Calculate estimated savings
        if stats['total_references'] > 0 and stats['unique_images'] > 0:
            avg_size = stats['unique_storage_bytes'] / stats['unique_images']
            estimated_without_dedup = stats['total_references'] * avg_size
            stats['estimated_saved_bytes'] = int(estimated_without_dedup - stats['unique_storage_bytes'])
            stats['deduplication_ratio'] = round(stats['total_references'] / stats['unique_images'], 2)
        else:
            stats['estimated_saved_bytes'] = 0
            stats['deduplication_ratio'] = 1.0

        return stats

    def get_document(self, document_id: int) -> Optional[Dict]:
        """Get document with content."""
        doc = self.db.get_document(document_id)
        if not doc:
            return None

        # Add file info if available
        if doc.get('local_path') and Path(doc['local_path']).exists():
            doc['file_exists'] = True
            doc['file_size'] = Path(doc['local_path']).stat().st_size
        else:
            doc['file_exists'] = False

        return doc

    def get_documents(
        self,
        meeting_id: int = None,
        agenda_item_id: int = None,
        search: str = None,
        limit: int = 50,
        offset: int = 0
    ) -> List[Dict]:
        """Get documents with filters."""
        return self.db.get_documents(
            meeting_id=meeting_id,
            agenda_item_id=agenda_item_id,
            search=search,
            limit=limit,
            offset=offset
        )

    def search_documents(self, query: str, limit: int = 20) -> List[Dict]:
        """
        Search documents by title and content.

        Args:
            query: Search query
            limit: Maximum results

        Returns:
            List of matching documents with relevance info
        """
        # Simple keyword search in database
        docs = self.db.get_documents(search=query, limit=limit)

        # Add match info
        query_lower = query.lower()
        for doc in docs:
            title_match = query_lower in (doc.get('title') or '').lower()
            content_match = query_lower in (doc.get('text_content') or '').lower()
            doc['match_type'] = []
            if title_match:
                doc['match_type'].append('title')
            if content_match:
                doc['match_type'].append('content')

        return docs

    def get_document_content(self, document_id: int) -> Optional[str]:
        """Get extracted text content of a document."""
        doc = self.db.get_document(document_id)
        if doc:
            return doc.get('text_content')
        return None

    def get_storage_stats(self) -> Dict:
        """Get storage statistics."""
        total_files = 0
        total_size = 0

        for file in self.documents_dir.iterdir():
            if file.is_file():
                total_files += 1
                total_size += file.stat().st_size

        db_stats = self.db.get_statistics()

        return {
            'files_on_disk': total_files,
            'total_size_bytes': total_size,
            'total_size_mb': round(total_size / (1024 * 1024), 2),
            'documents_in_db': db_stats.get('documents', 0),
            'documents_by_status': db_stats.get('documents_by_status', {})
        }


# Singleton instance
_provider_instance = None


def get_document_provider() -> DocumentProvider:
    """Get singleton document provider instance."""
    global _provider_instance
    if _provider_instance is None:
        _provider_instance = DocumentProvider()
    return _provider_instance


if __name__ == '__main__':
    # Test the provider
    provider = get_document_provider()

    print("Storage stats:")
    stats = provider.get_storage_stats()
    for key, value in stats.items():
        print(f"  {key}: {value}")

    print("\nPending documents:")
    from core.database import get_database
    db = get_database()
    pending = db.get_documents_pending_download()
    print(f"  {len(pending)} documents pending download")

    if pending[:3]:
        print("\nDownloading first 3 documents...")
        success, failed = provider.download_pending_documents(limit=3)
        print(f"  Downloaded: {success}, Failed: {failed}")
